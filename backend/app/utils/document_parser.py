"""
文档解析工具
支持PDF、PPT、Word文档解析，输出统一文本格式

使用库：
- pdfplumber: PDF解析
- python-docx: Word文档解析
- python-pptx: PPT解析
"""
import os
from typing import Optional, List, Tuple, Dict, Any
from pathlib import Path
import re

try:
    import pdfplumber
    from docx import Document
    from pptx import Presentation
except ImportError as e:
    print(f"警告: 文档解析库未安装，请运行 pip install pdfplumber python-docx python-pptx")
    print(f"错误详情: {e}")


class DocumentParser:
    """文档解析器基类，提供统一的文本格式输出"""
    
    @staticmethod
    def clean_text(text: str) -> str:
        """
        清理文本，统一格式
        
        Args:
            text: 原始文本
        
        Returns:
            清理后的文本
        """
        if not text:
            return ""
        
        # 移除多余的空白字符
        text = re.sub(r'\s+', ' ', text)
        # 移除多余的换行符（保留段落分隔）
        text = re.sub(r'\n{3,}', '\n\n', text)
        # 移除首尾空白
        text = text.strip()
        
        return text
    
    @staticmethod
    def format_output(content_parts: List[str], page_count: int, doc_type: str) -> str:
        """
        格式化输出文本
        
        Args:
            content_parts: 内容片段列表
            page_count: 页数/幻灯片数
            doc_type: 文档类型
        
        Returns:
            格式化后的统一文本
        """
        # 过滤空内容
        content_parts = [part.strip() for part in content_parts if part.strip()]
        
        if not content_parts:
            return ""
        
        # 合并内容
        full_text = "\n\n".join(content_parts)
        
        # 清理文本
        full_text = DocumentParser.clean_text(full_text)
        
        return full_text


def parse_pdf(file_path: str) -> Tuple[str, int]:
    """
    使用pdfplumber解析PDF文件
    
    Args:
        file_path: PDF文件路径
    
    Returns:
        (文本内容, 页数)
    
    Raises:
        Exception: 解析失败时抛出异常
    """
    try:
        content_parts = []
        page_count = 0
        
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                # 提取文本
                page_text = page.extract_text()
                
                if page_text:
                    # 清理页面文本
                    page_text = DocumentParser.clean_text(page_text)
                    if page_text:
                        # 添加页面标记
                        content_parts.append(f"=== 第 {page_num} 页 ===\n{page_text}")
                
                # 尝试提取表格
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, 1):
                        table_text = format_table(table)
                        if table_text:
                            content_parts.append(f"--- 第 {page_num} 页 表格 {table_num} ---\n{table_text}")
        
        # 格式化输出
        text_content = DocumentParser.format_output(content_parts, page_count, 'pdf')
        
        return text_content, page_count
    
    except Exception as e:
        raise Exception(f"PDF解析失败: {str(e)}")


def format_table(table: List[List]) -> str:
    """
    格式化表格为文本
    
    Args:
        table: 表格数据（二维列表）
    
    Returns:
        格式化后的表格文本
    """
    if not table:
        return ""
    
    rows = []
    for row in table:
        # 过滤None值并转换为字符串
        row_text = [str(cell) if cell is not None else "" for cell in row]
        rows.append(" | ".join(row_text))
    
    return "\n".join(rows)


def parse_word(file_path: str) -> Tuple[str, int]:
    """
    使用python-docx解析Word文档
    
    Args:
        file_path: Word文档路径
    
    Returns:
        (文本内容, 页数)
    
    Raises:
        Exception: 解析失败时抛出异常
    """
    try:
        content_parts = []
        
        doc = Document(file_path)
        
        # 解析段落
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if para_text:
                content_parts.append(para_text)
        
        # 解析表格
        for table_num, table in enumerate(doc.tables, 1):
            table_rows = []
            for row in table.rows:
                row_cells = [cell.text.strip() if cell.text else "" for cell in row.cells]
                row_text = " | ".join(row_cells)
                if row_text.strip():
                    table_rows.append(row_text)
            
            if table_rows:
                table_text = f"--- 表格 {table_num} ---\n" + "\n".join(table_rows)
                content_parts.append(table_text)
        
        # 估算页数（假设每页约500字）
        full_text = "\n".join(content_parts)
        page_count = max(1, len(full_text) // 500)
        
        # 格式化输出
        text_content = DocumentParser.format_output(content_parts, page_count, 'word')
        
        return text_content, page_count
    
    except Exception as e:
        raise Exception(f"Word文档解析失败: {str(e)}")


def parse_ppt(file_path: str) -> Tuple[str, int]:
    """
    使用python-pptx解析PPT文件
    
    Args:
        file_path: PPT文件路径
    
    Returns:
        (文本内容, 幻灯片数)
    
    Raises:
        Exception: 解析失败时抛出异常
    """
    try:
        content_parts = []
        
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 文本框
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_texts.append(text)
                
                # 表格
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() if cell.text else "" for cell in row.cells]
                        row_text = " | ".join(row_cells)
                        if row_text.strip():
                            table_rows.append(row_text)
                    
                    if table_rows:
                        table_text = "表格:\n" + "\n".join(table_rows)
                        slide_texts.append(table_text)
            
            # 如果有内容，添加到结果中
            if slide_texts:
                slide_content = "\n".join(slide_texts)
                slide_content = DocumentParser.clean_text(slide_content)
                if slide_content:
                    content_parts.append(f"=== 幻灯片 {slide_num} ===\n{slide_content}")
        
        # 格式化输出
        text_content = DocumentParser.format_output(content_parts, slide_count, 'ppt')
        
        return text_content, slide_count
    
    except Exception as e:
        raise Exception(f"PPT解析失败: {str(e)}")


def parse_document(file_path: str) -> Tuple[str, int, str]:
    """
    解析文档（统一入口）
    
    根据文件扩展名自动选择对应的解析器
    
    Args:
        file_path: 文档文件路径
    
    Returns:
        (文本内容, 页数/幻灯片数, 文件类型)
    
    Raises:
        ValueError: 不支持的文件类型
        Exception: 解析失败
    
    Examples:
        >>> text, pages, doc_type = parse_document("example.pdf")
        >>> print(f"提取了 {pages} 页内容")
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        text, page_count = parse_pdf(file_path)
        return text, page_count, 'pdf'
    
    elif file_ext in ['.doc', '.docx']:
        text, page_count = parse_word(file_path)
        return text, page_count, 'word'
    
    elif file_ext in ['.ppt', '.pptx']:
        text, page_count = parse_ppt(file_path)
        return text, page_count, 'ppt'
    
    else:
        raise ValueError(f"不支持的文件类型: {file_ext}。支持的类型: .pdf, .doc, .docx, .ppt, .pptx")


def get_document_info(file_path: str) -> Dict[str, Any]:
    """
    获取文档基本信息
    
    Args:
        file_path: 文档文件路径
    
    Returns:
        包含文档信息的字典
    """
    file_path_obj = Path(file_path)
    
    info = {
        'filename': file_path_obj.name,
        'file_type': file_path_obj.suffix.lower(),
        'file_size': os.path.getsize(file_path) if os.path.exists(file_path) else 0,
        'exists': os.path.exists(file_path)
    }
    
    # 如果文件存在，尝试获取页数
    if info['exists']:
        try:
            _, page_count, doc_type = parse_document(file_path)
            info['page_count'] = page_count
            info['document_type'] = doc_type
        except Exception as e:
            info['error'] = str(e)
    
    return info

